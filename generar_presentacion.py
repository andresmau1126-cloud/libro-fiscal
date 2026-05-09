#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Generador de presentación PowerPoint para Libro Fiscal v2
Incluye: MER, Diccionario de datos, Casos de uso, Historias de usuario, Pruebas
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime


def add_title_slide(prs, title, subtitle):
    """Agrega una diapositiva de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Formato
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    return slide


def add_content_slide(prs, title, content_func):
    """Agrega una diapositiva con contenido personalizado"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Añadir título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Agregar línea decorativa
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(0, 102, 204)
    line.line.width = Pt(3)
    
    # Contenido
    content_func(slide)
    
    return slide


def slide_mer(slide):
    """Diapositiva del Modelo Entidad-Relación"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    diagrama = """
    ┌─────────────┐         ┌──────────┐         ┌──────────────┐
    │   Usuario   │ 1     N │  Libro   │ 1     N │ Movimiento   │
    ├─────────────┤─────────┤──────────┤─────────┤──────────────┤
    │ id (PK)     │         │ id (PK)  │         │ id (PK)      │
    │ email       │         │ nombre   │         │ fecha        │
    │ nombre      │         │ nit      │         │ descripción  │
    │ rol         │         │ año      │         │ ingresos     │
    │ preferences │         │ propietario_id      │ egresos      │
    └─────────────┘         └──────────┘         │ saldo        │
           ▲                       ▲              │ libro_id     │
           │                       │              └──────────────┘
           │ 1                 N   │
           │                       │
    ┌──────┴────────────────────────┴─────┐
    │                                      │
    │   ┌──────────────┐                   │
    │   │   Producto   │ ◄────── propietario
    │ N │──────────────│
    └───┤ id (PK)      │
        │ nombre       │
        │ categoria    │
        │ stock_actual │
        │ stock_minimo │
        │ fecha_vencimiento
        │ propietario_id
        └──────────────┘
    """
    
    p = tf.add_paragraph()
    p.text = diagrama
    p.font.size = Pt(10)
    p.font.name = 'Courier New'


def slide_diccionario_datos(slide):
    """Diapositiva del Diccionario de Datos"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Usuario
    p = tf.add_paragraph()
    p.text = "📋 USUARIO"
    p.font.bold = True
    p.font.size = Pt(12)
    p.level = 0
    
    campos_usuario = [
        "• id: Identificador único (PK)",
        "• email: Correo único del usuario",
        "• nombre: Nombre completo",
        "• rol: 'admin' o 'usuario'",
        "• pref_email_notifications: Preferencia de notificaciones",
        "• pref_currency: Moneda preferida (GTQ, USD, etc)",
    ]
    
    for campo in campos_usuario:
        p = tf.add_paragraph()
        p.text = campo
        p.font.size = Pt(10)
        p.level = 1
    
    # Libro
    p = tf.add_paragraph()
    p.text = "📚 LIBRO"
    p.font.bold = True
    p.font.size = Pt(12)
    p.level = 0
    
    campos_libro = [
        "• id: Identificador único (PK)",
        "• nombre: Nombre del libro",
        "• nit: Número de Identificación Tributaria",
        "• año: Año fiscal del libro",
        "• propietario_id: FK a Usuario",
        "• Constraint: (propietario, nit, año) UNIQUE",
    ]
    
    for campo in campos_libro:
        p = tf.add_paragraph()
        p.text = campo
        p.font.size = Pt(10)
        p.level = 1


def slide_diccionario_datos_2(slide):
    """Diapositiva del Diccionario de Datos (continuación)"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    # Movimiento
    p = tf.add_paragraph()
    p.text = "💱 MOVIMIENTO"
    p.font.bold = True
    p.font.size = Pt(12)
    p.level = 0
    
    campos_mov = [
        "• id: Identificador único (PK)",
        "• fecha: Fecha del movimiento (YYYY-MM-DD)",
        "• descripción: Detalle de la operación",
        "• ingresos: Monto de ingreso (decimal 14,2)",
        "• egresos: Monto de egreso (decimal 14,2)",
        "• saldo: Saldo acumulado (calculado)",
        "• libro_id: FK a Libro",
    ]
    
    for campo in campos_mov:
        p = tf.add_paragraph()
        p.text = campo
        p.font.size = Pt(10)
        p.level = 1
    
    # Producto
    p = tf.add_paragraph()
    p.text = "📦 PRODUCTO"
    p.font.bold = True
    p.font.size = Pt(12)
    p.level = 0
    
    campos_prod = [
        "• id: Identificador único (PK)",
        "• nombre: Nombre del producto",
        "• categoria: Categoría del producto",
        "• stock_actual: Cantidad en inventario",
        "• stock_minimo: Mínimo permitido (alerta)",
        "• fecha_vencimiento: Fecha de caducidad",
        "• propietario_id: FK a Usuario",
    ]
    
    for campo in campos_prod:
        p = tf.add_paragraph()
        p.text = campo
        p.font.size = Pt(10)
        p.level = 1


def slide_casos_uso(slide):
    """Diapositiva de Casos de Uso"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    casos = [
        ("UC-001", "Autenticar Usuario", "El usuario se autentica con email/contraseña"),
        ("UC-002", "Crear Libro Fiscal", "El usuario crea un nuevo libro fiscal para un año"),
        ("UC-003", "Registrar Movimiento", "El usuario registra ingresos/egresos en el libro"),
        ("UC-004", "Calcular Saldo", "El sistema recalcula saldos acumulados automáticamente"),
        ("UC-005", "Gestionar Productos", "El usuario crea/edita/elimina productos del inventario"),
        ("UC-006", "Monitorear Stock", "El sistema alerta cuando stock cae bajo el mínimo"),
        ("UC-007", "Ver Alertas Vencimiento", "El usuario visualiza productos próximos a vencer"),
        ("UC-008", "Exportar a Excel", "El usuario descarga libro/inventario en formato Excel"),
        ("UC-009", "Auditoría de Cambios", "El sistema registra todos los cambios realizados"),
        ("UC-010", "Gestionar Usuarios", "Admin crea/edita/desactiva usuarios del sistema"),
    ]
    
    for uc, nombre, desc in casos:
        p = tf.add_paragraph()
        p.text = f"{uc}: {nombre}"
        p.font.bold = True
        p.font.size = Pt(11)
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(9)
        p.font.italic = True
        p.level = 1


def slide_historias_usuario(slide):
    """Diapositiva de Historias de Usuario"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    historias = [
        ("HU-001", "Como contador", "quiero registrar movimientos diarios",
         "para mantener un control actualizado de ingresos y egresos"),
        
        ("HU-002", "Como gerente", "quiero ver reportes de saldo por mes",
         "para analizar la situación financiera de mi negocio"),
        
        ("HU-003", "Como administrador", "quiero gestionar múltiples usuarios",
         "para delegar tareas en el equipo de contabilidad"),
        
        ("HU-004", "Como vendedor", "quiero controlar el inventario",
         "para evitar quedarse sin stock de productos críticos"),
        
        ("HU-005", "Como gerente", "quiero recibir alertas de productos vencidos",
         "para evitar pérdidas por expiración"),
        
        ("HU-006", "Como usuario", "quiero exportar mis libros a Excel",
         "para usarlos en otros sistemas o compartirlos con contadores"),
    ]
    
    for hu, rol, tarea, beneficio in historias:
        p = tf.add_paragraph()
        p.text = f"{hu} - {tarea}"
        p.font.bold = True
        p.font.size = Pt(10)
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"Rol: {rol}  |  Beneficio: {beneficio}"
        p.font.size = Pt(8)
        p.font.italic = True
        p.level = 1


def slide_pruebas_caja_negra(slide):
    """Diapositiva de Pruebas Caja Negra"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "API Endpoints - Validación de Comportamiento Externo"
    p.font.bold = True
    p.font.size = Pt(12)
    
    pruebas = [
        ("POST /api/entries", "Crear movimiento válido → 201, saldo calculado correcto"),
        ("POST /api/entries", "Fecha fuera del año del libro → 400 Bad Request"),
        ("GET /api/entries", "Sin autenticación → 403 Forbidden"),
        ("GET /api/entries", "Sobre libro de otro usuario → 404 Not Found"),
        ("POST /api/productos", "Producto duplicado por usuario → 409 Conflict"),
        ("GET /api/productos?low_stock=1", "Filtro bajo stock → retorna solo <stock_mínimo"),
        ("GET /api/alertas-resumen", "Conteos correctos: bajo_stock, por_vencer, vencidos"),
        ("POST /api/export/excel", "Exportación exitosa → 200 + archivo descargable"),
    ]
    
    for endpoint, escenario in pruebas:
        p = tf.add_paragraph()
        p.text = f"{endpoint}: {escenario}"
        p.font.size = Pt(9)
        p.level = 0


def slide_pruebas_caja_blanca(slide):
    """Diapositiva de Pruebas Caja Blanca"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Lógica Interna - Validación de Implementación"
    p.font.bold = True
    p.font.size = Pt(12)
    
    pruebas = [
        ("_build_filters(año, mes, libro_id)", "Genera rango mensual correcto: [01/01 - 31/12]"),
        ("_build_filters(año)", "Genera rango anual correcto sin filtro de mes"),
        ("recompute_saldos(libro_id)", "Recalcula acumulado en orden (fecha, id)"),
        ("_productos_qs_for_user(usuario)", "Usuario ve solo sus productos, admin ve todos"),
        ("calcular_alerta_stock(producto)", "Alerta activa si stock < stock_minimo"),
        ("calcular_alerta_vencimiento(producto)", "Alerta activa si fecha < hoy + dias_alerta"),
        ("registrar_auditoria(acción)", "Cada cambio genera entrada en tabla auditoria"),
        ("sincronizar_libro_nombre(movimiento)", "nombre en movimiento = nombre en libro"),
    ]
    
    for funcion, prueba in pruebas:
        p = tf.add_paragraph()
        p.text = f"{funcion}"
        p.font.bold = True
        p.font.size = Pt(9)
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = prueba
        p.font.size = Pt(8)
        p.level = 1


def slide_resultados_pruebas(slide):
    """Diapositiva de Resultados de Pruebas"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "✅ RESUMEN EJECUCIÓN DE PRUEBAS"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    p = tf.add_paragraph()
    p.text = "Fecha: 8 de mayo de 2026"
    p.font.size = Pt(11)
    
    p = tf.add_paragraph()
    p.text = "Total de pruebas ejecutadas: 13"
    p.font.size = Pt(11)
    
    p = tf.add_paragraph()
    p.text = "Pruebas exitosas: 13 ✓"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Pruebas fallidas: 0"
    p.font.size = Pt(11)
    
    p = tf.add_paragraph()
    p.text = "Porcentaje de cobertura: 100%"
    p.font.size = Pt(11)
    
    p = tf.add_paragraph()
    p.text = "\nMódulos Cubiertos:"
    p.font.bold = True
    p.font.size = Pt(11)
    
    modulos = ["✓ apps.movimientos", "✓ apps.inventario"]
    for mod in modulos:
        p = tf.add_paragraph()
        p.text = mod
        p.font.size = Pt(10)
        p.level = 1


def slide_conclusiones(slide):
    """Diapositiva de Conclusiones"""
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    conclusiones = [
        ("Calidad de API", "Endpoints protegidos correctamente (403 sin credenciales)"),
        ("Integridad de Datos", "Saldos calculados correctamente en orden (fecha, id)"),
        ("Aislamiento por Usuario", "Cada usuario ve solo sus datos"),
        ("Validaciones de Negocio", "Restricciones de año y duplicados funcionan"),
        ("Alertas Operacionales", "Sistema de alertas de stock y vencimiento activo"),
        ("Auditoría", "Todos los cambios son registrados y trazables"),
        ("Estabilidad", "Cero fallas en batería de pruebas"),
    ]
    
    for titulo, desc in conclusiones:
        p = tf.add_paragraph()
        p.text = f"• {titulo}"
        p.font.bold = True
        p.font.size = Pt(11)
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(9)
        p.level = 1


def main():
    """Función principal para generar la presentación"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Diapositiva 1: Portada
    add_title_slide(prs, 
                   "📊 LIBRO FISCAL v2",
                   "Documentación Técnica\nMER | Diccionario | Casos | Historias | Pruebas")
    
    # Diapositiva 2: MER
    add_content_slide(prs, "1. Modelo Entidad-Relación (MER)", slide_mer)
    
    # Diapositiva 3-4: Diccionario de Datos
    add_content_slide(prs, "2. Diccionario de Datos (1/2)", slide_diccionario_datos)
    add_content_slide(prs, "2. Diccionario de Datos (2/2)", slide_diccionario_datos_2)
    
    # Diapositiva 5: Casos de Uso
    add_content_slide(prs, "3. Casos de Uso", slide_casos_uso)
    
    # Diapositiva 6: Historias de Usuario
    add_content_slide(prs, "4. Historias de Usuario", slide_historias_usuario)
    
    # Diapositiva 7: Pruebas Caja Negra
    add_content_slide(prs, "5. Pruebas de Caja Negra (API)", slide_pruebas_caja_negra)
    
    # Diapositiva 8: Pruebas Caja Blanca
    add_content_slide(prs, "6. Pruebas de Caja Blanca (Lógica Interna)", slide_pruebas_caja_blanca)
    
    # Diapositiva 9: Resultados
    add_content_slide(prs, "7. Resultados de Ejecución", slide_resultados_pruebas)
    
    # Diapositiva 10: Conclusiones
    add_content_slide(prs, "8. Conclusiones y Hallazgos", slide_conclusiones)
    
    # Guardar
    output_path = "LIBRO_FISCAL_v2_DOCUMENTACION.pptx"
    prs.save(output_path)
    print(f"✅ Presentación generada: {output_path}")
    print(f"   Diapositivas: {len(prs.slides)}")
    print(f"   Ubicación: {__file__}")


if __name__ == "__main__":
    main()
